import pandas as pd
from Toolkit import Position, add_employees_to_slide
import pptx
import datetime
from data_preparation import data_preparation_processor
from constants import TEST_PICS_PATH, OUTPUT_DATA_PATH

def generate_presentation_from_template(prs: pptx.presentation.Presentation, raw_df: pd.DataFrame) -> None:
    # Get separate dataframes for different teams and job families

    teams_df, job_families_df , all_data_df = data_preparation_processor(raw_df)

    business_research_df, office_management_df, business_translation_df, it_df, finance_df, \
        graphic_design_df, marketing_df, operations_df, data_analytics_df, business_development_df, \
        human_capital_df, executive_management_df, sales_representative_df = job_families_df

    mck_team_df, office_admin_team_df, ik_team_df, bcg_me_team_df, business_translation_team_df, \
        it_team_df, finance_team_df, bcg_europe_team_df, graphic_design_team_df, bcg_africa_team_df, \
        marketing_team_df, bcg_kt_support_df, operations_team_df, data_analytics_services_df, \
        business_development_team_df, hr_team_df, research_team_df, executive_management_df, ikt_finance_df = teams_df

    # retrieve all the job families in seperate dfs
    ########## The Cover SLide : Slide[0] , displays the presentation title : 'WHo's WHo ' and the date in ' Month/Year Format #########################################
    # add the current date to the slide
    # get slide 0 of the presentation
    slide = prs.slides[0]
    # get the first textbox shape in the slide
    textbox1 = slide.shapes[0]
    # create a new textbox shape
    textbox2 = slide.shapes.add_textbox(textbox1.left, textbox1.top + textbox1.height + 20, textbox1.width, 50)
    # add the current month and year to the textbox
    today = datetime.date.today()
    month_year_str = today.strftime('%B %Y')
    textbox2.text_frame.text = month_year_str

    # format the text in the new textbox
    font = textbox2.text_frame.paragraphs[0].font
    font.size = pptx.util.Pt(30)  # set font size to 20 points
    font.bold = True  # set bold text
    font.color.rgb = pptx.dml.color.RGBColor(255, 255, 255)  # set white text color

    ############ SLIDE [1] : HARD CODED : Managment Commitee ###############################
    #### This slide is directly derived from the input presentation : it remains unchaged till a further notice ###########

    ######################### SLIDE [2] : HARD CODED :  #######################################################
    #### This slide is directly derived from the input presentation : it remains unchaged till a further notice #################

    ######################### SLIDE [3] : BD Team  #######################################################
    # create dfs to be used in slide 3

    slide = prs.slides[3]
    bd_officer_df = business_development_df[(business_development_df['Is_officer'] == 1)]
    bd_sr_vp_df = business_development_df[
        (business_development_df['Is_vice_president'] == 1) & (business_development_df['Is_senior'] == 1)]
    bd_vp_df = business_development_df[
        (business_development_df['Is_vice_president'] == 1) & (business_development_df['Is_senior'] == 0) & (
                    business_development_df['Is_associate'] == 0)]
    bd_associate_vps_df = business_development_df[
        (business_development_df['Is_associate'] == 1) & (business_development_df['Is_vice_president'] == 1)]
    bd_lead_df = business_development_df[(business_development_df['Is_lead'] == 1)]
    bd_sr_acc_exec_df = business_development_df[
        (business_development_df['Is_acc_exec'] == 1) & (business_development_df['Is_senior'] == 1)]
    bd_acc_exec_df = business_development_df[
        (business_development_df['Is_acc_exec'] == 1) & (business_development_df['Is_senior'] == 0)]


    add_employees_to_slide(df=bd_officer_df, row_length=2, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_MIDDLE_1 , need_flag=1)
    add_employees_to_slide(df=bd_sr_vp_df, row_length=3, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=0, position=Position.TOP_LEFT_4 , need_flag=1)
    add_employees_to_slide(df=bd_vp_df, row_length=4, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=0, position=Position.TOP_MIDDLE_4 , need_flag=1)
    add_employees_to_slide(df=bd_associate_vps_df, row_length=8, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_1,need_flag=1)
    add_employees_to_slide(df=bd_lead_df, row_length=4, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=0, position=Position.BOTTOM_MIDDLE_1,need_flag=1)
    add_employees_to_slide(df=bd_sr_acc_exec_df, row_length=5, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_3,need_flag=1)
    add_employees_to_slide(df=bd_acc_exec_df, row_length=5, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=0, position=Position.BOTTOM_MIDDLE_3,need_flag=1)

    ######################### SLIDE [4] : Finance Team  #######################################################

    finance_directors_df = finance_team_df[finance_team_df['Is_director'] == 1]
    finance_senior_coordinators_df = finance_team_df[
        (finance_team_df['Is_senior'] == 1) & (finance_team_df['Is_coordinator'] == 1)]
    finance_others_df = finance_team_df[~((finance_team_df['Is_director'] == 1) | (
            (finance_team_df['Is_senior'] == 1) & (finance_team_df['Is_coordinator'] == 1)))]

    slide = prs.slides[4]

    # Finance Strategic Topics
    add_employees_to_slide(df=finance_directors_df, row_length=2, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_MIDDLE_3 , need_flag=1)
    # Business Controller Topics, supporting BD team
    add_employees_to_slide(df=finance_senior_coordinators_df, row_length=3, font_size=7, slide=slide,
                           pics_path=TEST_PICS_PATH, detailed_slide=1, position=Position.BOTTOM_MIDDLE_2 , need_flag=1)

    # Country Finance Topics
    add_employees_to_slide(df=finance_others_df, row_length=3, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_RIGHT_2 , need_flag=1)


    ######################### SLIDE [5] : HR Team   #######################################################
    ########## This SLide is left blank : pending detailed Data ###################################

    HR_executive_df = hr_team_df[(hr_team_df['Is_Executive_Management'] == 1)]
    HR_recruitment_df = hr_team_df[(hr_team_df['Is_Recruitment'] == 1)]
    HR_PD_df = hr_team_df[(hr_team_df['Is_PD'] == 1)]
    HR_Administration_df = hr_team_df[(hr_team_df['Is_HR_Administration'] == 1)]

    slide = prs.slides[5]

    add_employees_to_slide(df=HR_executive_df, row_length=2, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_MIDDLE_3 , need_flag=1)

    add_employees_to_slide(df=HR_recruitment_df, row_length=3, font_size=7, slide=slide,
                           pics_path=TEST_PICS_PATH, detailed_slide=1, position=Position.BOTTOM_MIDDLE_2 , need_flag=1)


    add_employees_to_slide(df=HR_PD_df, row_length=2, font_size=6, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_RIGHT_2 , need_flag=1)


    add_employees_to_slide(df=HR_Administration_df, row_length=3, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.BOTTOM_RIGHT_1 , need_flag=1)


    ######################### SLIDE [6] : Operations Team  #######################################################

    COO_df = operations_team_df[operations_team_df['Is_COO'] == 1]

    # partners_df = operations_team_df[operations_team_df['Is_partner'] == 1]

    operations_team_country_managers_df = operations_team_df[
        (operations_team_df['Is_COO'] != 1) & (operations_team_df['Is_partner'] != 1)
        & (operations_team_df['Is_country_manager'] == 1)]

    operations_team_others_df = operations_team_df[
        (operations_team_df['Is_COO'] != 1) & (operations_team_df['Is_partner'] != 1)
        & (operations_team_df['Is_country_manager'] != 1)]

    slide = prs.slides[6]
    add_employees_to_slide(df=COO_df, row_length=1, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1,
                           position=Position.MIDDLE_LEFT_3 , need_flag=1)
    add_employees_to_slide(df=operations_team_country_managers_df, row_length=8, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.MAIN_MIDDLE_LEFT , need_flag=1)
    add_employees_to_slide(df=operations_team_others_df, row_length=4, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.BOTTOM_LEFT_3 , need_flag=1)

    ######################### SLIDE [7] :   IT Team #######################################################
    it_team_managers_df = it_team_df[(it_team_df['Is_manager'] == 1)]
    it_team_others_df = it_team_df[(it_team_df['Is_manager'] != 1)]


    slide = prs.slides[7]
    add_employees_to_slide(df=it_team_managers_df, row_length=4, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_3 , need_flag=1)
    add_employees_to_slide(df=it_team_others_df, row_length=8, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.BOTTOM_LEFT_1 , need_flag=1)

    ######################### SLIDE [8] :   Marketing Team #######################################################
    slide = prs.slides[8]
    add_employees_to_slide(df=marketing_team_df, row_length=4, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_3 , need_flag=1)

    ######################### SLIDE [9] : Office Admin Team  #######################################################


    slide = prs.slides[9]
    add_employees_to_slide(df=office_admin_team_df, row_length=3, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_3 , need_flag=1)

    ######################### SLIDE [10] : Infomineo Research Teams Presentation #######################################################
    ####### HARD CODED ######################################

    ######################### SLIDE [11] : Business Reserach Leaders  #######################################################
    # create dfs to be used in slide 10
    # Create a dataframe of only the research team head(s)
    research_team_df = business_research_df
    #research_head_df = research_team_df[research_team_df['Is_Head'] == 1]
    # Create a dataframe of only the research managers
    research_leaders_df = research_team_df[research_team_df['Is_BR_leader'] == 1]
    #research_team_df = business_research_df

    slide = prs.slides[11]

    add_employees_to_slide(df=research_leaders_df, row_length=13, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_2 , need_flag=1)

    ######################### SLIDE [12] : Business Research 2/4 : IKT Team  #######################################################
    # create dfs to be used in slide 11
    # Create a dataframe of only the research team head(s)

    ikt_members_egypt_df = ik_team_df[
        (ik_team_df['Is_BR_leader'] == 0) & ik_team_df['location'].str.contains('egypt', case=False)]

    ikt_members_mexico_df = ik_team_df[
        (ik_team_df['Is_BR_leader'] == 0) & ik_team_df['location'].str.contains('mexico', case=False)]

    slide = prs.slides[12]
    add_employees_to_slide(df=ikt_members_egypt_df, row_length=13, font_size=6, slide=slide,
                           pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_2 , need_flag=0)
    add_employees_to_slide(df=ikt_members_mexico_df, row_length=12, font_size=7, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.BOTTOM_LEFT_3 , need_flag=0)

    ######################### SLIDE [13] :  IKT Morocco / Remote   #######################################################

    # Create a dataframe of only the senior research senior analysts

    ikt_members_morocco_df = ik_team_df[
        (ik_team_df['Is_BR_leader'] == 0) & ik_team_df['location'].str.contains('morocco', case=False)]

    ikt_members_remote_df = ik_team_df[
        (ik_team_df['Is_BR_leader'] == 0) & ik_team_df['location'].str.contains('remote', case=False)]

    slide = prs.slides[13]

    add_employees_to_slide(df=ikt_members_morocco_df, row_length=12, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_2,need_flag=0)
    add_employees_to_slide(df=ikt_members_remote_df, row_length=12, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.BOTTOM_LEFT_3 , need_flag=0)


    ######################### SLIDE [14] :  MCK Team members    #######################################################
    # create dfs to be used in slide 16
    # Create a dataframe of only the research analysts
    mck_members_egypt_df = mck_team_df[
        (mck_team_df['Is_BR_leader'] == 0) & mck_team_df['location'].str.contains('egypt', case=False)]

    mck_members_mexico_df = mck_team_df[
        (mck_team_df['Is_BR_leader'] == 0) & mck_team_df['location'].str.contains('mexico', case=False)]

    #research_analyst_df = research_team_df[(research_team_df['Is_analyst'] == 1) & (research_team_df['Is_senior'] == 0)]

    slide = prs.slides[14]

    add_employees_to_slide(df=mck_members_egypt_df, row_length=12, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_2,need_flag=0)
    add_employees_to_slide(df=mck_members_mexico_df, row_length=12, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.BOTTOM_LEFT_2 ,need_flag=0)

    ######################### SLIDE [15] : MCK Team 2/2  #######################################################
    # Hard coded : Left as is for the time being
    # Create a dataframe of only the research analysts

    mck_members_morocco_df = mck_team_df[
        (mck_team_df['Is_BR_leader'] == 0) & mck_team_df['location'].str.contains('morocco', case=False)]

    mck_members_remote_df = mck_team_df[
        (mck_team_df['Is_BR_leader'] == 0) & mck_team_df['location'].str.contains('remote', case=False)]

    slide = prs.slides[15]

    add_employees_to_slide(df=mck_members_morocco_df, row_length=12, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_2,need_flag=0)
    add_employees_to_slide(df=mck_members_remote_df, row_length=12, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.BOTTOM_LEFT_3,need_flag=0)


    ######################### SLIDE [16] : BCG members 1    #######################################################
    # create dfs to be used in slide 17
    # Create a dataframe of only the research analysts
    bcg_team_df = pd.concat([bcg_me_team_df, bcg_europe_team_df, bcg_africa_team_df, bcg_kt_support_df])
    # Filter employees from Egypt in BCG team
    bcg_members_egypt_df = bcg_team_df[
        (bcg_team_df['Is_BR_leader'] == 0) & bcg_team_df['location'].str.contains('egypt', case=False)
        ]

    # Filter employees from Mexico in BCG team
    bcg_members_mexico_df = bcg_team_df[
        (bcg_team_df['Is_BR_leader'] == 0) & bcg_team_df['location'].str.contains('mexico', case=False)
        ]

    # Assuming you have the necessary functions like add_employees_to_slide and Position enumeration defined:
    # Slide for adding employees
    slide = prs.slides[16]

    # Add employees from Egypt to the slide
    add_employees_to_slide(df=bcg_members_egypt_df, row_length=12, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_2,need_flag=0)

    # Add employees from Mexico to the slide
    add_employees_to_slide(df=bcg_members_mexico_df, row_length=12, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.BOTTOM_LEFT_2,need_flag=0)


    ######################### SLIDE [17] : BCG Team 2/2  #######################################################
    # create dfs to be used in slide 18
    # Create a dataframe of only the research analysts
    # Filter employees from Morocco in BCG team
    bcg_members_morocco_df = bcg_team_df[
        (bcg_team_df['Is_BR_leader'] == 0) & bcg_team_df['location'].str.contains('morocco', case=False)
        ]

    # Filter employees from Remote location in BCG team
    bcg_members_remote_df = bcg_team_df[
        (bcg_team_df['Is_BR_leader'] == 0) & bcg_team_df['location'].str.contains('remote', case=False)
        ]

    # Assuming you have the necessary functions like add_employees_to_slide and Position enumeration defined:
    # Slide for adding employees
    slide = prs.slides[17]

    # Add employees from Morocco to the slide
    add_employees_to_slide(df=bcg_members_morocco_df, row_length=12, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.TOP_LEFT_2,need_flag=0)

    # Add employees from Remote location to the slide
    add_employees_to_slide(df=bcg_members_remote_df, row_length=12, font_size=8, slide=slide, pics_path=TEST_PICS_PATH,
                           detailed_slide=1, position=Position.BOTTOM_LEFT_3,need_flag=0)


    ######################### SLIDE [18] :  Business Translation  #######################################################
    # create dfs to be used in slide 18
    # Create a dataframe of only the research analysts
    business_translation_manager_df = business_translation_df[business_translation_df['Is_BR_leader'] == 1]
    business_translation_others_df = business_translation_df[business_translation_df['Is_BR_leader'] == 0]
    slide = prs.slides[18]
    add_employees_to_slide(df=business_translation_manager_df, row_length=5, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH, detailed_slide=1, position=Position.TOP_LEFT_2,need_flag=1)

    add_employees_to_slide(df=business_translation_others_df, row_length=12, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH, detailed_slide=1, position=Position.TOP_LEFT_5,need_flag=1)




    ######################### SLIDE [19] : Graphic Design  #######################################################
    # create dfs to be used in slide 20
    # Filter business translation manager and others in graphic design team
    graphic_design_manager_df = graphic_design_team_df[graphic_design_team_df['Is_BR_leader'] == 1]
    graphic_design_others_df = graphic_design_team_df[graphic_design_team_df['Is_BR_leader'] == 0]

    # Assuming you have the necessary functions like add_employees_to_slide and Position enumeration defined:
    # Slide for adding employees
    slide = prs.slides[19]


    # Add business translation managers to the slide
    add_employees_to_slide(df=graphic_design_manager_df, row_length=5, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH, detailed_slide=1, position=Position.TOP_LEFT_2,need_flag=1)

    # Add other employees from graphic design team to the slide
    add_employees_to_slide(df=graphic_design_others_df, row_length=12, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH, detailed_slide=1, position=Position.TOP_LEFT_5,need_flag=1)

    ######################### SLIDE [20] : BCG Team : EUR / ME / Africa / Pilot (2/2)  #######################################################
    # create dfs to be used in slide 22
    # Create a dataframe of only the research analysts
    # Filter data analytics services manager and others
    data_analytics_manager_df = data_analytics_services_df[data_analytics_services_df['Is_BR_leader'] == 1]
    data_analytics_others_df = data_analytics_services_df[data_analytics_services_df['Is_BR_leader'] == 0]

    # Assuming you have the necessary functions like add_employees_to_slide and Position enumeration defined:
    # Slide for adding employees
    slide = prs.slides[20]

    # Add data analytics services managers to the slide
    add_employees_to_slide(df=data_analytics_manager_df, row_length=2, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH, detailed_slide=1, position=Position.TOP_MIDDLE_1,need_flag=1)

    # Add other employees from data analytics services team to the slide
    add_employees_to_slide(df=data_analytics_others_df, row_length=8, font_size=8, slide=slide,
                           pics_path=TEST_PICS_PATH, detailed_slide=1, position=Position.MIDDLE_LEFT_4,need_flag=1)


    return prs


